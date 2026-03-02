"""
评测报告生成器
将评测结果生成 Markdown、JSON、Word 和 PPT 格式的报告
"""

import os
import json
import traceback
from datetime import datetime
from typing import Dict, Any, List
from pathlib import Path


class ReportGenerator:
    """评测报告生成器"""

    def __init__(self, output_dir: str = None):
        """
        初始化报告生成器

        Args:
            output_dir: 输出目录
        """
        if output_dir is None:
            output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "outputs")

        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def generate(self, evaluation_result: Dict[str, Any], formats: List[str] = None) -> str:
        """
        生成评测报告

        Args:
            evaluation_result: 评测结果
            formats: 输出格式列表 ["markdown", "json"]

        Returns:
            生成的文件路径
        """
        import logging
        logger = logging.getLogger(__name__)

        if formats is None:
            formats = ["markdown"]

        script_name = evaluation_result.get("script_name", "unnamed")
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        logger.info(f"开始生成报告: {script_name}, 格式: {formats}")

        generated_files = []

        if "markdown" in formats:
            logger.info("生成 Markdown 报告...")
            md_path = self._generate_markdown(evaluation_result, script_name, timestamp)
            generated_files.append(md_path)
            logger.info(f"Markdown 报告已生成: {md_path}")

        if "json" in formats:
            logger.info("生成 JSON 报告...")
            json_path = self._generate_json(evaluation_result, script_name, timestamp)
            generated_files.append(json_path)
            logger.info(f"JSON 报告已生成: {json_path}")

        if "word" in formats or "docx" in formats:
            logger.info("生成 Word 报告...")
            word_path = self._generate_word(evaluation_result, script_name, timestamp)
            generated_files.append(word_path)
            logger.info(f"Word 报告已生成: {word_path}")

        if "ppt" in formats or "pptx" in formats:
            logger.info("生成 PPT 报告...")
            ppt_path = self._generate_ppt(evaluation_result, script_name, timestamp)
            generated_files.append(ppt_path)
            logger.info(f"PPT 报告已生成: {ppt_path}")

        logger.info(f"所有报告生成完成: {len(generated_files)} 个文件")
        return generated_files

    def _generate_markdown(
        self,
        result: Dict[str, Any],
        script_name: str,
        timestamp: str
    ) -> str:
        """
        生成 Markdown 格式报告

        Args:
            result: 评测结果
            script_name: 剧本名称
            timestamp: 时间戳

        Returns:
            生成的文件路径
        """
        import logging
        logger = logging.getLogger(__name__)

        filename = f"{script_name}_{timestamp}.md"
        filepath = os.path.join(self.output_dir, filename)

        logger.info(f"准备生成 Markdown 报告: {filename}")

        try:
            with open(filepath, 'w', encoding='utf-8') as f:
                # 标题
                overall = result.get("overall", {})
                f.write(f"# 《{script_name}》剧本评测报告\n\n")
                f.write(f"> 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")

                # 综合评分
                f.write("## 综合评分\n\n")
                score = overall.get("total_score", 0)
                grade = overall.get("grade", "N/A")
                f.write(f"### {score}/100  |  等级: **{grade}**\n\n")

                # 分项评分表
                f.write("## 分项评分\n\n")
                f.write("| 维度 | 得分 | 满分 | 权重 | 加权得分 |\n")
                f.write("|------|------|------|------|----------|\n")

                for detail in overall.get("details", []):
                    dim = detail.get("dimension", "")
                    score = detail.get("score", 0)
                    max_score = detail.get("max_score", 100)
                    weight = detail.get("weight", 0)
                    weighted = detail.get("weighted_score", 0)
                    f.write(f"| {dim} | {score} | {max_score} | {weight*100:.0f}% | {weighted:.2f} |\n")

                f.write("\n")

                # 各维度详细分析
                f.write("## 详细分析\n\n")

                for dim_key, dim_result in result.get("dimensions", {}).items():
                    if "error" in dim_result:
                        f.write(f"### {dim_result.get('dimension_name', dim_key)}\n\n")
                        f.write(f"❌ 评测失败: {dim_result['error']}\n\n")
                        f.write("---\n\n")
                        continue

                    f.write(f"### {dim_result.get('dimension_name', dim_key)}\n\n")
                    f.write(f"**得分**: {dim_result.get('total_score', 0)}/{dim_result.get('max_score', 100)}\n\n")

                    # 子项得分
                    sub_scores = dim_result.get('sub_scores', {})
                    if sub_scores:
                        f.write("#### 📊 子项评分\n\n")
                        f.write("| 项目 | 得分 | 满分 | 评价 |\n")
                        f.write("|------|------|------|------|\n")
                        for sub_key, sub_value in sub_scores.items():
                            name = sub_value.get("name", sub_key)
                            score = sub_value.get("score", 0)
                            max_score = sub_value.get("max_score", 100)
                            comment = sub_value.get("comment", "")
                            # 截断过长的评论
                            if len(comment) > 100:
                                comment = comment[:97] + "..."
                            f.write(f"| {name} | {score} | {max_score} | {comment} |\n")
                        f.write("\n")

                    # 扣分项
                    penalties = dim_result.get('penalties', [])
                    penalties_applied = dim_result.get('penalties_applied', [])
                    all_penalties = penalties if penalties else penalties_applied

                    if all_penalties:
                        f.write("#### ❌ 扣分项\n\n")
                        total_penalty = 0
                        for penalty in all_penalties:
                            item = penalty.get('item', penalty.get('name', '未知项'))
                            score = penalty.get('score', 0)
                            reason = penalty.get('reason', penalty.get('details', ''))
                            total_penalty += abs(score) if score < 0 else 0
                            f.write(f"- **{item}**: {score}分")
                            if reason:
                                f.write(f"\n  - *{reason}*")
                            f.write("\n")
                        f.write(f"\n**累计扣分**: {total_penalty:.0f}分\n\n")

                    # 优点
                    strengths = dim_result.get('strengths', [])
                    if strengths:
                        f.write("#### ✅ 优点\n\n")
                        for i, strength in enumerate(strengths, 1):
                            f.write(f"{i}. {strength}\n")
                        f.write("\n")

                    # 待改进点
                    weaknesses = dim_result.get('weaknesses', [])
                    if weaknesses:
                        f.write("#### ⚠️ 待改进点\n\n")
                        for i, weakness in enumerate(weaknesses, 1):
                            f.write(f"{i}. {weakness}\n")
                        f.write("\n")

                    # 改进建议
                    suggestions = dim_result.get('suggestions', [])
                    if suggestions:
                        f.write("#### 💡 改进建议\n\n")
                        for i, suggestion in enumerate(suggestions, 1):
                            f.write(f"**建议 {i}**: {suggestion}\n\n")
                        f.write("\n")

                    # 特殊内容
                    if "notable_lines" in dim_result:
                        f.write("#### 精彩台词\n\n")
                        for line in dim_result["notable_lines"]:
                            if isinstance(line, dict):
                                f.write(f"> **{line.get('speaker', '')}**: {line.get('line', '')}\n")
                                f.write(f"> \n> *{line.get('reason', '')}*\n\n")
                            elif isinstance(line, str):
                                f.write(f"> {line}\n\n")

                    if "character_analysis" in dim_result:
                        f.write("#### 人物分析\n\n")
                        for char in dim_result["character_analysis"]:
                            if isinstance(char, dict):
                                f.write(f"**{char.get('character', '')}** ({char.get('role', '')}) - "
                                       f"{char.get('score', 0)}/{char.get('max_score', 10)}\n\n")
                                f.write(f"{char.get('analysis', '')}\n\n")
                            elif isinstance(char, str):
                                f.write(f"{char}\n\n")

                    if "twists_identified" in dim_result:
                        f.write("#### 反转分析\n\n")
                        for twist in dim_result["twists_identified"]:
                            f.write(f"- **{twist.get('position', '')}**: {twist.get('description', '')} "
                                   f"({twist.get('effectiveness_score', 0)}/{twist.get('max_score', 10)})\n")
                        f.write("\n")

                    if "target_audience" in dim_result:
                        audience = dim_result["target_audience"]
                        f.write("#### 目标受众\n\n")
                        if audience.get("primary"):
                            f.write(f"- **主要受众**: {', '.join(audience['primary'])}\n")
                        if audience.get("age_range"):
                            f.write(f"- **年龄范围**: {audience['age_range']}\n")
                        if audience.get("gender_preference"):
                            f.write(f"- **性别偏好**: {audience['gender_preference']}\n")
                        if audience.get("interest_tags"):
                            f.write(f"- **兴趣标签**: {', '.join(audience['interest_tags'])}\n")
                        f.write("\n")

                    # 每个维度后添加分隔线
                    f.write("---\n\n")

                # 总结建议
                f.write("## 📋 总结建议\n\n")

                f.write("### 🌟 核心优势\n\n")
                all_strengths = []
                for dim_result in result.get("dimensions", {}).values():
                    dim_name = dim_result.get('dimension_name', '')
                    dim_strengths = dim_result.get("strengths", [])
                    for strength in dim_strengths:
                        all_strengths.append(f"[{dim_name}] {strength}")

                # 显示所有优点，不限制数量
                for i, strength in enumerate(all_strengths, 1):
                    f.write(f"{i}. {strength}\n")
                f.write("\n")

                f.write("### 🔧 重点改进方向\n\n")
                all_suggestions = []
                for dim_result in result.get("dimensions", {}).values():
                    dim_name = dim_result.get('dimension_name', '')
                    dim_suggestions = dim_result.get("suggestions", [])
                    for suggestion in dim_suggestions:
                        all_suggestions.append(f"[{dim_name}] {suggestion}")

                # 显示所有建议
                for i, suggestion in enumerate(all_suggestions, 1):
                    f.write(f"{i}. {suggestion}\n")
                f.write("\n")

                f.write("### 📈 综合评价\n\n")
                overall = result.get("overall", {})
                total_score = overall.get("total_score", 0)
                grade = overall.get("grade", "N/A")

                if grade == 'A' or grade == 'S':
                    f.write(f"🎉 恭喜！您的剧本获得了 **{grade}** 级评价（{total_score}分），属于优秀水平。\n\n")
                    f.write("剧本展现了出色的创作能力，各方面表现均衡且突出。建议保持当前水准，并在细节上继续打磨。\n\n")
                elif grade == 'B':
                    f.write(f"👍 您的剧本获得了 **{grade}** 级评价（{total_score}分），属于良好水平。\n\n")
                    f.write("剧本整体表现良好，具备一定竞争力。建议根据上述改进方向进行优化，有望提升到更高等级。\n\n")
                elif grade == 'C':
                    f.write(f"💪 您的剧本获得了 **{grade}** 级评价（{total_score}分），尚有改进空间。\n\n")
                    f.write("建议重点关注上述待改进点和改进建议，进行系统性修改，以提升剧本质量和市场竞争力。\n\n")
                else:
                    f.write(f"📝 您的剧本获得了 **{grade}** 级评价（{total_score}分），建议进行大幅修改。\n\n")
                    f.write("建议从故事结构、人物塑造、对话质量等多个维度进行全面优化，参考上述改进建议逐项改进。\n\n")

                f.write("---\n")
                f.write("*本报告由 AI 剧本评测系统基于豆包 seed-1.8 模型生成，仅供参考。如需更精准的分析，建议结合专业人工评审。*\n")

                logger.info(f"Markdown 报告写入成功: {filepath}")
        except Exception as e:
            logger.error(f"Markdown 报告生成失败: {e}")
            import traceback
            logger.error(traceback.format_exc())
            raise

        return filepath

    def _generate_json(
        self,
        result: Dict[str, Any],
        script_name: str,
        timestamp: str
    ) -> str:
        """
        生成 JSON 格式报告

        Args:
            result: 评测结果
            script_name: 剧本名称
            timestamp: 时间戳

        Returns:
            生成的文件路径
        """
        import logging
        logger = logging.getLogger(__name__)

        filename = f"{script_name}_{timestamp}.json"
        filepath = os.path.join(self.output_dir, filename)

        logger.info(f"准备生成 JSON 报告: {filename}")

        # 添加元数据
        result["metadata"] = {
            "generated_at": datetime.now().isoformat(),
            "generator": "AI Script Evaluator v1.0",
            "script_name": script_name
        }

        try:
            logger.info("开始写入 JSON 文件...")
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False, indent=2)
            logger.info(f"JSON 报告写入成功: {filepath}")
        except Exception as e:
            logger.error(f"JSON 报告生成失败: {e}")
            import traceback
            logger.error(traceback.format_exc())
            raise

        return filepath

    def generate_batch_summary(
        self,
        results: List[Dict[str, Any]]
    ) -> str:
        """
        生成批量评测汇总报告

        Args:
            results: 评测结果列表

        Returns:
            生成的文件路径
        """
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"batch_summary_{timestamp}.md"
        filepath = os.path.join(self.output_dir, filename)

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write("# 批量评测汇总报告\n\n")
            f.write(f"> 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"> 评测剧本数量: {len(results)}\n\n")

            # 排行榜
            f.write("## 评测排行榜\n\n")
            f.write("| 排名 | 剧本名称 | 综合得分 | 等级 |\n")
            f.write("|------|----------|----------|------|\n")

            sorted_results = sorted(
                results,
                key=lambda x: x.get("overall", {}).get("total_score", 0),
                reverse=True
            )

            for i, result in enumerate(sorted_results, 1):
                script_name = result.get("script_name", "Unknown")
                score = result.get("overall", {}).get("total_score", 0)
                grade = result.get("overall", {}).get("grade", "N/A")
                f.write(f"| {i} | {script_name} | {score} | {grade} |\n")

            f.write("\n")

            # 统计信息
            f.write("## 统计信息\n\n")
            scores = [r.get("overall", {}).get("total_score", 0) for r in results]
            if scores:
                avg_score = sum(scores) / len(scores)
                max_score = max(scores)
                min_score = min(scores)
                f.write(f"- **平均分**: {avg_score:.2f}\n")
                f.write(f"- **最高分**: {max_score} ({sorted_results[0].get('script_name', 'Unknown')})\n")
                f.write(f"- **最低分**: {min_score}\n\n")

            # 等级分布
            grade_count = {}
            for result in results:
                grade = result.get("overall", {}).get("grade", "N/A")
                grade_count[grade] = grade_count.get(grade, 0) + 1

            f.write("### 等级分布\n\n")
            for grade in ["S", "A", "B", "C", "D"]:
                count = grade_count.get(grade, 0)
                bar = "█" * count
                f.write(f"- **{grade}级**: {count} {bar}\n")
            f.write("\n")

        return filepath

    def _generate_word(
        self,
        result: Dict[str, Any],
        script_name: str,
        timestamp: str
    ) -> str:
        """
        生成 Word 格式报告

        Args:
            result: 评测结果
            script_name: 剧本名称
            timestamp: 时间戳

        Returns:
            生成的文件路径
        """
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

        doc = Document()
        
        # 标题
        title = doc.add_heading(f"《{script_name}》剧本评测报告", 0)
        title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        # 生成时间
        doc.add_paragraph(f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

        # 综合评分
        overall = result.get("overall", {})
        total_score = overall.get("total_score", 0)
        grade = overall.get("grade", "N/A")
        
        doc.add_heading("综合评分", 1)
        p = doc.add_paragraph()
        run = p.add_run(f"{total_score}/100")
        run.font.size = Pt(32)
        run.font.color.rgb = RGBColor(102, 126, 234)
        run.bold = True
        
        doc.add_paragraph(f"等级: {grade}")

        # 分项评分
        doc.add_heading("分项评分", 1)
        table = doc.add_table(rows=1, cols=4)
        table.style = 'Light Grid Accent 1'
        
        # 表头
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = "维度"
        hdr_cells[1].text = "得分"
        hdr_cells[2].text = "满分"
        hdr_cells[3].text = "权重"

        # 数据行
        dimensions = result.get("dimensions", {})
        for dim_key, dim_result in dimensions.items():
            if "error" not in dim_result:
                row_cells = table.add_row().cells
                row_cells[0].text = dim_result.get("dimension_name", dim_key)
                row_cells[1].text = str(dim_result.get("total_score", 0))
                row_cells[2].text = str(dim_result.get("max_score", 100))
                # 权重从配置中获取，这里简化处理
                row_cells[3].text = "5%"

        # 详细分析
        doc.add_page_break()
        doc.add_heading("详细分析", 1)

        for dim_key, dim_result in dimensions.items():
            if "error" in dim_result:
                continue
            
            doc.add_heading(dim_result.get("dimension_name", dim_key), 2)
            doc.add_paragraph(f"得分: {dim_result.get('total_score', 0)}/{dim_result.get('max_score', 100)}")

            # 子项评分
            if "sub_scores" in dim_result:
                doc.add_paragraph("子项评分:", style="List Bullet")
                for sub_key, sub_data in dim_result["sub_scores"].items():
                    p = doc.add_paragraph()
                    p.add_run(f"{sub_data.get('name', sub_key)}: ").bold = True
                    p.add_run(f"{sub_data.get('score', 0)}/{sub_data.get('max_score', 100)}")
                    doc.add_paragraph(f"评价: {sub_data.get('comment', '')}")

            # 扣分项
            penalties = dim_result.get('penalties', [])
            penalties_applied = dim_result.get('penalties_applied', [])
            all_penalties = penalties if penalties else penalties_applied

            if all_penalties:
                doc.add_paragraph("扣分项:", style="List Bullet")
                total_penalty = 0
                for penalty in all_penalties:
                    item = penalty.get('item', penalty.get('name', '未知项'))
                    score = penalty.get('score', 0)
                    reason = penalty.get('reason', penalty.get('details', ''))
                    total_penalty += abs(score) if score < 0 else 0
                    p = doc.add_paragraph(style="List Bullet 2")
                    p.add_run(f"{item}: {score}分").bold = True
                    if reason:
                        doc.add_paragraph(f"原因: {reason}", style="List Bullet 3")
                p = doc.add_paragraph(style="List Bullet 2")
                p.add_run(f"累计扣分: {total_penalty:.0f}分").font.color.rgb = RGBColor(255, 0, 0)

            # 优点
            if "strengths" in dim_result and dim_result["strengths"]:
                doc.add_paragraph("优点:", style="List Bullet")
                for strength in dim_result["strengths"]:
                    doc.add_paragraph(strength, style="List Bullet 2")

            # 待改进点
            if "weaknesses" in dim_result and dim_result["weaknesses"]:
                doc.add_paragraph("待改进点:", style="List Bullet")
                for weakness in dim_result["weaknesses"]:
                    doc.add_paragraph(weakness, style="List Bullet 2")

            # 建议
            if "suggestions" in dim_result and dim_result["suggestions"]:
                doc.add_paragraph("改进建议:", style="List Bullet")
                for suggestion in dim_result["suggestions"]:
                    doc.add_paragraph(suggestion, style="List Bullet 2")

        # 保存文件
        filename = f"{script_name}_{timestamp}.docx"
        filepath = os.path.join(self.output_dir, filename)
        doc.save(filepath)

        return filepath

    def _generate_ppt(
        self,
        result: Dict[str, Any],
        script_name: str,
        timestamp: str
    ) -> str:
        """
        生成 PPT 格式报告

        Args:
            result: 评测结果
            script_name: 剧本名称
            timestamp: 时间戳

        Returns:
            生成的文件路径
        """
        from pptx import Presentation
        from pptx.util import Pt, Inches
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # 标题页
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        title.text = f"《{script_name}》剧本评测报告"
        
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

        # 综合评分页
        bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
        shapes = bullet_slide.shapes
        title_shape = shapes.title
        title_shape.text = "综合评分"

        overall = result.get("overall", {})
        total_score = overall.get("total_score", 0)
        grade = overall.get("grade", "N/A")

        body_shape = shapes.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.text = f"总分: {total_score}/100\n等级: {grade}"

        # 分项评分页
        dimensions = result.get("dimensions", {})
        
        # 每页显示3个维度
        dim_items = list(dimensions.items())
        for i in range(0, len(dim_items), 3):
            bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
            shapes = bullet_slide.shapes
            title_shape = shapes.title
            title_shape.text = f"分项评分 ({i//3 + 1})"

            body_shape = shapes.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()

            for j in range(3):
                if i + j < len(dim_items):
                    dim_key, dim_result = dim_items[i + j]
                    if "error" not in dim_result:
                        p = text_frame.add_paragraph()
                        p.level = j
                        p.text = f"{dim_result.get('dimension_name', dim_key)}: {dim_result.get('total_score', 0)}/{dim_result.get('max_score', 100)}"
                        p.font.size = Pt(18)

        # 详细分析页（每个维度一页）
        for dim_key, dim_result in dimensions.items():
            if "error" in dim_result:
                continue

            bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
            shapes = bullet_slide.shapes
            title_shape = shapes.title
            title_shape.text = dim_result.get("dimension_name", dim_key)

            body_shape = shapes.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()

            p = text_frame.add_paragraph()
            p.text = f"得分: {dim_result.get('total_score', 0)}/{dim_result.get('max_score', 100)}"
            p.font.size = Pt(24)
            p.font.bold = True

            # 扣分项
            penalties = dim_result.get('penalties', [])
            penalties_applied = dim_result.get('penalties_applied', [])
            all_penalties = penalties if penalties else penalties_applied

            if all_penalties:
                p = text_frame.add_paragraph()
                p.text = "扣分项:"
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 0, 0)
                p.level = 1
                total_penalty = 0
                for penalty in all_penalties[:5]:  # 最多显示5个扣分项
                    item = penalty.get('item', penalty.get('name', '未知项'))
                    score = penalty.get('score', 0)
                    reason = penalty.get('reason', penalty.get('details', ''))
                    total_penalty += abs(score) if score < 0 else 0
                    p = text_frame.add_paragraph()
                    p.text = f"• {item}: {score}分"
                    p.level = 2
                    if reason and len(reason) < 50:  # 只显示较短的原因
                        p = text_frame.add_paragraph()
                        p.text = f"  {reason}"
                        p.level = 3
                # 显示累计扣分
                p = text_frame.add_paragraph()
                p.text = f"累计扣分: {total_penalty:.0f}分"
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 0, 0)
                p.level = 2

            # 优点
            if "strengths" in dim_result and dim_result["strengths"]:
                p = text_frame.add_paragraph()
                p.text = "优点:"
                p.font.bold = True
                p.level = 1
                for strength in dim_result["strengths"][:3]:
                    p = text_frame.add_paragraph()
                    p.text = f"• {strength}"
                    p.level = 2

            # 待改进点
            if "weaknesses" in dim_result and dim_result["weaknesses"]:
                p = text_frame.add_paragraph()
                p.text = "待改进点:"
                p.font.bold = True
                p.level = 1
                for weakness in dim_result["weaknesses"][:3]:
                    p = text_frame.add_paragraph()
                    p.text = f"• {weakness}"
                    p.level = 2

            # 建议
            if "suggestions" in dim_result and dim_result["suggestions"]:
                p = text_frame.add_paragraph()
                p.text = "改进建议:"
                p.font.bold = True
                p.level = 1
                for suggestion in dim_result["suggestions"][:3]:
                    p = text_frame.add_paragraph()
                    p.text = f"• {suggestion}"
                    p.level = 2

        # 保存文件
        filename = f"{script_name}_{timestamp}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)

        return filepath
